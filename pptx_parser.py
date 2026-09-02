# -*- coding: utf-8 -*-
"""
PPTX 내 모든 슬라이드의 모든 표를 순회해 공통 스키마(list of dict)로 변환.
공통 스키마:
{
  "slide": int, "table_index": int, "row": int, "column": int,
  "field": str|None, "product": str|None, "value": str
}
병합 셀은 origin 셀에서만 값을 채우고, spanned(병합에 흡수된) 셀은 스킵해서
중복 값이 생기지 않게 한다. field/product 매핑은 여기서 하지 않고
raw 그리드(행렬)만 만들어 반환 -> 이후 화면에서 사용자가 방향을 확인(7장)한 뒤 매핑한다.
"""
from pptx import Presentation
from ..utils import clean_whitespace


def parse_pptx_tables(file_obj) -> list:
    """
    반환: [
      {
        "slide": int, "table_index": int,
        "grid": [[cell_text, ...], ...],   # 병합 셀은 origin 위치에만 값, spanned는 "" (occupied 표시는 merge_map에)
        "merge_map": [[bool, ...], ...],   # True면 이 셀이 다른 셀에 병합되어 흡수된 셀(빈칸 취급)
      }, ...
    ]
    """
    prs = Presentation(file_obj)
    results = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        table_idx_in_slide = 0
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table_idx_in_slide += 1
            table = shape.table
            n_rows = len(table.rows)
            n_cols = len(table.columns)

            grid = [["" for _ in range(n_cols)] for _ in range(n_rows)]
            merge_map = [[False for _ in range(n_cols)] for _ in range(n_rows)]

            for r in range(n_rows):
                for c in range(n_cols):
                    cell = table.cell(r, c)
                    # python-pptx: is_merge_origin=True인 셀만 실제 텍스트를 채택.
                    # is_spanned=True(다른 셀에 병합되어 흡수됨)면 빈 칸으로 두고 merge_map만 표시.
                    if getattr(cell, "is_spanned", False) and not getattr(cell, "is_merge_origin", False):
                        merge_map[r][c] = True
                        continue
                    text = clean_whitespace(cell.text)
                    grid[r][c] = text

            results.append({
                "slide": slide_idx,
                "table_index": table_idx_in_slide,
                "grid": grid,
                "merge_map": merge_map,
            })

    return results
